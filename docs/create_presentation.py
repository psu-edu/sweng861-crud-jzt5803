"""
Campus Analytics Platform – SWENG 861 Presentation Generator
Generates a professional 9-slide PowerPoint presentation.
Run: python3 docs/create_presentation.py

Assertion-Evidence format (Michael Alley, Penn State):
  - Slide TITLE  = one complete-sentence assertion (the claim)
  - Slide BODY   = minimal evidence: diagram, table, or short key points
  - Speaker VOICE = carries the detail and explanation
  - Rule: slide must be readable in ~5 seconds
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
NAVY    = RGBColor(0x1e, 0x3a, 0x5f)
STEEL   = RGBColor(0x21, 0x96, 0xf3)
LBLUE   = RGBColor(0x90, 0xca, 0xf9)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LGRAY   = RGBColor(0xf0, 0xf4, 0xf8)
DGRAY   = RGBColor(0x33, 0x33, 0x33)
MONO_BG = RGBColor(0xf5, 0xf5, 0xf5)
RED_LT  = RGBColor(0xff, 0xeb, 0xee)
RED_DK  = RGBColor(0xc6, 0x28, 0x28)
GRN_LT  = RGBColor(0xe8, 0xf5, 0xe9)
GRN_DK  = RGBColor(0x2e, 0x7d, 0x32)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUTPUT_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "campus-analytics-presentation.pptx"
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _cell_bg(cell, rgb):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    h = '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2])
    sf = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{h}"/></a:solidFill>'
    )
    for e in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(e)
    tcPr.insert(0, sf)


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def textbox(slide, text, l, t, w, h,
            size=14, bold=False, color=DGRAY,
            align=PP_ALIGN.LEFT, italic=False,
            wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def multiline_box(slide, lines_data, l, t, w, h,
                  default_size=13, bg=None, mono=False):
    """
    lines_data: list of (text, size, bold, color, indent_level)
    Each tuple renders as one paragraph.
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
    else:
        tb.fill.background()

    for i, item in enumerate(lines_data):
        text = item[0]
        size = item[1] if len(item) > 1 else default_size
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else DGRAY
        level = item[4] if len(item) > 4 else 0

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Courier New" if mono else "Calibri"
    return tb


def slide_header(slide, assertion_text):
    """Navy bar + complete-sentence assertion as slide title."""
    rect(slide, Inches(0),    Inches(0),    SLIDE_W,       Inches(0.08), NAVY)
    rect(slide, Inches(0),    Inches(0.08), Inches(0.06),  Inches(1.0),  STEEL)
    rect(slide, Inches(0.06), Inches(0.08), Inches(13.27), Inches(1.0),  LGRAY)
    textbox(slide, assertion_text,
            Inches(0.28), Inches(0.12), Inches(12.8), Inches(0.92),
            size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT, wrap=False)
    rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), NAVY)


def section_label(slide, text, l, t, w, h=Inches(0.32), color=NAVY):
    """Small bold column / section label."""
    textbox(slide, text, l, t, w, h, size=12, bold=True, color=color, wrap=False)


def divider(slide, x):
    """Vertical steel divider at position x."""
    rect(slide, x, Inches(1.1), Inches(0.04), Inches(6.3), STEEL)


def add_table(slide, headers, rows, l, t, w, h, col_widths=None):
    """Navy-header table."""
    nr = len(rows) + 1
    nc = len(headers)
    tbl = slide.shapes.add_table(nr, nc, l, t, w, h).table

    if col_widths:
        for i, cw in enumerate(col_widths[:nc]):
            tbl.columns[i].width = cw

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        _cell_bg(cell, NAVY)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.size = Pt(11)
        r.font.name = "Calibri"

    for ri, row in enumerate(rows):
        bg = LGRAY if ri % 2 == 0 else WHITE
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = txt
            _cell_bg(cell, bg)
            p = cell.text_frame.paragraphs[0]
            r = p.runs[0]
            r.font.size = Pt(10)
            r.font.name = "Calibri"
            r.font.color.rgb = DGRAY
            if ci == 0:
                r.font.bold = True
                r.font.color.rgb = NAVY
    return tbl


# ---------------------------------------------------------------------------
# Slide 1 – Title (no assertion format needed)
# ---------------------------------------------------------------------------

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    rect(slide, Inches(0), Inches(3.1), SLIDE_W, Inches(0.05), STEEL)

    textbox(slide, "Campus Analytics Platform",
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
            size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, "SWENG 861 – Software Construction  |  Spring 2026",
            Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.8),
            size=22, bold=False, color=LBLUE, align=PP_ALIGN.CENTER)
    textbox(slide, "Jomar Thomas Almonte",
            Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.65),
            size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, "Pennsylvania State University",
            Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.55),
            size=15, color=LBLUE, align=PP_ALIGN.CENTER)
    return slide


# ---------------------------------------------------------------------------
# Slide 2 – Requirements
# Assertion: Campus Analytics Eliminates University Metrics Data Silos
# Evidence: The 3 failures  vs  7 delivered features
# ---------------------------------------------------------------------------

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "Campus Analytics Eliminates University Metrics Data Silos")

    # ── LEFT: The Problem ──────────────────────────────────────────────────
    section_label(slide, "The Problem", Inches(0.3), Inches(1.15), Inches(5.9))

    problem = [
        ("Departments track metrics in spreadsheets.", 13, False, DGRAY),
        ("That creates three specific failures:", 13, False, DGRAY),
        ("", 6),
        ("✗  No audit trail — who changed what, and when?", 13, False, RED_DK),
        ("✗  No access control — any user can edit any data", 13, False, RED_DK),
        ("✗  No single source of truth — data silos per dept", 13, False, RED_DK),
        ("", 6),
        ("Who is affected:", 12, True, NAVY),
        ("Faculty          enrollment trends & academic KPIs", 12, False, DGRAY),
        ("Dept heads       facilities & financial monitoring", 12, False, DGRAY),
        ("Admins           cross-domain oversight & audit", 12, False, DGRAY),
    ]
    multiline_box(slide, problem, Inches(0.3), Inches(1.5), Inches(5.9), Inches(5.8))

    divider(slide, Inches(6.35))

    # ── RIGHT: The Solution ────────────────────────────────────────────────
    section_label(slide, "The Solution: 7 Features Shipped", Inches(6.55), Inches(1.15), Inches(6.5))

    solution = [
        ("✓  JWT + OAuth auth, bcrypt, rate limiting", 13, False, GRN_DK),
        ("✓  5-category metrics CRUD + admin master view", 13, False, GRN_DK),
        ("✓  Live weather widget (°F / mph, Penn State)", 13, False, GRN_DK),
        ("✓  Domain event audit trail — append-only", 13, False, GRN_DK),
        ("✓  Prometheus + Grafana observability", 13, False, GRN_DK),
        ("✓  GitHub Actions CI/CD pipeline", 13, False, GRN_DK),
        ("✓  Non-root Docker, 3-stage multi-stage build", 13, False, GRN_DK),
        ("", 6),
        ("Non-Functional Requirements — met, not aspirational:", 12, True, NAVY),
        ("Security     OWASP Top 10 mitigations applied", 12, False, DGRAY),
        ("Reliability  ≥99.5% availability, p95 ≤ 500 ms", 12, False, DGRAY),
        ("Testing      320 automated tests block broken builds", 12, False, DGRAY),
    ]
    multiline_box(slide, solution, Inches(6.55), Inches(1.5), Inches(6.5), Inches(5.8))
    return slide


# ---------------------------------------------------------------------------
# Slide 3 – Architecture
# Assertion: Next.js 15 Collapses Frontend and API Into One Deploy Unit
# Evidence: Architecture diagram (the diagram IS the proof)
# ---------------------------------------------------------------------------

def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "Next.js 15 Collapses Frontend and API Into One Deploy Unit")

    # ── LEFT: short rationale ──────────────────────────────────────────────
    left = [
        ("One codebase. One Docker image. No split deploy.", 13, True, NAVY),
        ("", 5),
        ("Frontend", 12, True, NAVY),
        ("Next.js 15 App Router (SSR) + React 19", 12, False, DGRAY),
        ("TailwindCSS 4 — zero custom CSS files", 12, False, DGRAY),
        ("", 5),
        ("API — 15 route handlers", 12, True, NAVY),
        ("JWT Bearer + NextAuth session auth", 12, False, DGRAY),
        ("4-tier rate limiting + BOLA owner checks", 12, False, DGRAY),
        ("", 5),
        ("Data — SQLite + Sequelize ORM", 12, True, NAVY),
        ("4 models: User, Metric, WeatherData, DomainEvent", 12, False, DGRAY),
        ("CASCADE deletes enforce referential integrity", 12, False, DGRAY),
        ("", 5),
        ("External", 12, True, NAVY),
        ("Open-Meteo weather  (free, keyless)", 12, False, DGRAY),
        ("Google OAuth 2.0 via NextAuth  (optional)", 12, False, DGRAY),
    ]
    multiline_box(slide, left, Inches(0.3), Inches(1.2), Inches(5.5), Inches(6.1))

    # ── RIGHT: architecture diagram ────────────────────────────────────────
    diagram = (
        "┌──────────────────────────────────┐\n"
        "│  Browser (React 19 + Next.js 15) │\n"
        "│  Dashboard │ Login │ Metrics      │\n"
        "│  Weather Widget │ Forms           │\n"
        "└─────────────────┬────────────────┘\n"
        "                  │ HTTPS / JWT\n"
        "┌─────────────────▼────────────────┐\n"
        "│  Next.js API Routes (Node.js)    │\n"
        "│  Auth │ Rate Limit │ Validation  │\n"
        "│  15 route handlers               │\n"
        "└──────┬──────────────────┬────────┘\n"
        "       │                  │\n"
        "┌──────▼──────┐  ┌────────▼───────┐\n"
        "│  SQLite DB   │  │ Open-Meteo API │\n"
        "│  Sequelize   │  │ (cached 10min) │\n"
        "│  4 models    │  └────────────────┘\n"
        "└──────┬──────┘\n"
        "       │\n"
        "┌──────▼──────────────────────────┐\n"
        "│  DevOps Layer                   │\n"
        "│  Docker │ GitHub Actions        │\n"
        "│  Prometheus │ Grafana           │\n"
        "└─────────────────────────────────┘"
    )

    tb = slide.shapes.add_textbox(Inches(6.1), Inches(1.2), Inches(6.9), Inches(6.1))
    tf = tb.text_frame
    tf.word_wrap = False
    tb.fill.solid()
    tb.fill.fore_color.rgb = MONO_BG
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = diagram
    run.font.size = Pt(9.5)
    run.font.name = "Courier New"
    run.font.color.rgb = DGRAY
    return slide


# ---------------------------------------------------------------------------
# Slide 4 – Tech Stack
# Assertion: Every Technology Choice Minimizes Overhead and Maximizes Testability
# Evidence: The choice table — each row shows the "why", not just the "what"
# ---------------------------------------------------------------------------

def build_slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "Every Technology Choice Minimizes Overhead and Maximizes Testability")

    headers = ["Layer", "Technology", "Why This Choice — The Evidence"]
    rows = [
        ["Frontend",       "Next.js 15 App Router",
         "UI + API in one image — no separate server, no CORS, no split deploy"],
        ["UI",             "React 19 + TailwindCSS 4",
         "Utility-first styling: zero custom CSS files in the entire codebase"],
        ["Auth",           "NextAuth 4 + JWT",
         "One getAuthUser() handles both token types — routes are unaware of the difference"],
        ["Database",       "SQLite + Sequelize 6",
         "File-based DB starts with zero infrastructure; ORM makes models unit-testable"],
        ["3rd Party",      "Open-Meteo Weather API",
         "Free and keyless — zero secrets to manage, rotate, or accidentally leak"],
        ["Testing",        "Jest 30 + React Testing Library",
         "~320 self-contained tests run in CI; broken builds never reach Docker"],
        ["Containers",     "Docker (3-stage multi-stage)",
         "Alpine + non-root user nextjs:1001 — minimal image, minimal attack surface"],
        ["CI/CD",          "GitHub Actions",
         "push → test → audit → Docker package — no manual deploy steps ever"],
        ["Observability",  "Prometheus + Grafana",
         "4 metrics, 5-panel SLO dashboard, zero external cloud dependency"],
    ]

    add_table(slide, headers, rows,
              Inches(0.25), Inches(1.2), Inches(12.83), Inches(6.05),
              col_widths=[Inches(1.35), Inches(2.3), Inches(9.18)])
    return slide


# ---------------------------------------------------------------------------
# Slide 5 – Core Features
# Assertion: Multi-Layer Defense and Seven Features Deliver Production-Grade Security
# Evidence: Two concise columns — auth/security left, data features right
# ---------------------------------------------------------------------------

def build_slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "Multi-Layer Defense and Seven Features Deliver Production-Grade Security")

    # ── LEFT: Auth & Security ──────────────────────────────────────────────
    section_label(slide, "Auth & Security", Inches(0.3), Inches(1.15), Inches(6.0))

    left = [
        ("Three token paths — one code path:", 13, True, NAVY),
        ("API clients       JWT Bearer (1-hr expiry)", 13, False, DGRAY),
        ("Browser users     NextAuth session cookie", 13, False, DGRAY),
        ("OAuth users       /api/auth/token bridge", 13, False, DGRAY),
        ("", 5),
        ("One getAuthUser() handles all three transparently.", 12, False, DGRAY),
        ("", 5),
        ("Brute force blocked at four layers:", 13, True, NAVY),
        ("Rate limit: 5 auth requests / 15 min / IP", 13, False, DGRAY),
        ("bcrypt cost 12 → deliberate 200 ms+ per hash", 13, False, DGRAY),
        ("JWT expiry limits stolen-token blast radius", 13, False, DGRAY),
        ("", 5),
        ("Five security headers on every response:", 13, True, NAVY),
        ("HSTS  |  X-Frame-Options: DENY  |  nosniff", 13, False, DGRAY),
        ("Owner BOLA check on all CRUD routes", 13, False, DGRAY),
        ("Sensitive keys stripped from all log output", 13, False, DGRAY),
    ]
    multiline_box(slide, left, Inches(0.3), Inches(1.52), Inches(6.0), Inches(5.8))

    divider(slide, Inches(6.45))

    # ── RIGHT: Data Features ───────────────────────────────────────────────
    section_label(slide, "Data Features", Inches(6.65), Inches(1.15), Inches(6.4))

    right = [
        ("Metrics CRUD — validated at every boundary:", 13, True, NAVY),
        ("5 enum categories (server-enforced)", 13, False, DGRAY),
        ("UUID keys prevent ID enumeration attacks", 13, False, DGRAY),
        ("Admin sees ALL; users see only their own", 13, False, DGRAY),
        ("US format: 47,892 students  |  $2.1B USD", 13, False, DGRAY),
        ("", 5),
        ("Weather — fetched once, cached ten minutes:", 13, True, NAVY),
        ("Open-Meteo, Penn State coords (40.7983°N)", 13, False, DGRAY),
        ("US units: °F temperature, mph wind speed", 13, False, DGRAY),
        ("", 5),
        ("Audit trail — every mutation, append-only:", 13, True, NAVY),
        ("metric.created / updated / deleted", 13, False, DGRAY),
        ("Async write — does not slow API response", 13, False, DGRAY),
        ("Events are never modified — true audit log", 13, False, DGRAY),
    ]
    multiline_box(slide, right, Inches(6.65), Inches(1.52), Inches(6.4), Inches(5.8))
    return slide


# ---------------------------------------------------------------------------
# Slide 6 – DevOps
# Assertion: Every Commit Is Automatically Built, Tested, and Packaged — No Manual Steps
# Evidence: Pipeline flow (left) + observability summary (right)
# ---------------------------------------------------------------------------

def build_slide6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "Every Commit Is Automatically Built, Tested, and Packaged — No Manual Steps")

    # ── LEFT: pipeline as ASCII flow ───────────────────────────────────────
    section_label(slide, "CI/CD Pipeline  (triggers on every push / PR to main)",
                  Inches(0.3), Inches(1.15), Inches(6.1))

    pipeline = (
        "push / PR to main\n"
        "       │\n"
        "  ┌────▼────────────────────────────┐\n"
        "  │  JOB 1 — build-and-test         │\n"
        "  │  npm ci  (lockfile-exact)        │\n"
        "  │  npm run build                   │\n"
        "  │  npm run test:coverage  ◄────────┼── QUALITY GATE\n"
        "  │    no continue-on-error          │   failure stops here\n"
        "  │  npm audit --audit-level=high    │\n"
        "  └────┬────────────────────────────┘\n"
        "       │  only if JOB 1 passes\n"
        "  ┌────▼────────────────────────────┐\n"
        "  │  JOB 2 — docker-build           │\n"
        "  │  3-stage build (Alpine)          │\n"
        "  │  non-root user nextjs:1001       │\n"
        "  │  smoke test: curl /api/health×5  │\n"
        "  └─────────────────────────────────┘"
    )

    tb = slide.shapes.add_textbox(Inches(0.3), Inches(1.52), Inches(6.1), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = False
    tb.fill.solid()
    tb.fill.fore_color.rgb = MONO_BG
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = pipeline
    r.font.size = Pt(10)
    r.font.name = "Courier New"
    r.font.color.rgb = DGRAY

    divider(slide, Inches(6.55))

    # ── RIGHT: observability ───────────────────────────────────────────────
    section_label(slide, "Observability Stack", Inches(6.75), Inches(1.15), Inches(6.3))

    right = [
        ("Three health endpoints — each has one role:", 13, True, NAVY),
        ("/api/health        load balancer heartbeat (always 200)", 12, False, DGRAY),
        ("/api/health/live   Kubernetes liveness probe (always 200)", 12, False, DGRAY),
        ("/api/health/ready  readiness probe (503 if DB is down)", 12, False, DGRAY),
        ("", 5),
        ("Structured JSON logging:", 13, True, NAVY),
        ("{timestamp, level, message, …meta}", 12, False, DGRAY),
        ("password / token / secret — stripped on every write", 12, False, DGRAY),
        ("", 5),
        ("Four Prometheus metric families:", 13, True, NAVY),
        ("http_requests_total          traffic by route + status", 12, False, DGRAY),
        ("http_request_duration_ms     latency histogram", 12, False, DGRAY),
        ("metrics_created_total        domain KPI counter", 12, False, DGRAY),
        ("auth_logins_total            security event signal", 12, False, DGRAY),
        ("", 5),
        ("Grafana: 5-panel dashboard, auto-provisioned on start", 12, False, DGRAY),
        ("SLOs: availability ≥ 99.5%  |  p95 latency ≤ 500 ms", 13, True, NAVY),
    ]
    multiline_box(slide, right, Inches(6.75), Inches(1.52), Inches(6.3), Inches(5.8))
    return slide


# ---------------------------------------------------------------------------
# Slide 7 – AI Audit
# Assertion: AI-Generated Code Contained 3 Security Flaws — All Caught by Manual Review
# Evidence: Flaw table (the audit results) + challenge list
# ---------------------------------------------------------------------------

def build_slide7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "AI-Generated Code Contained 3 Security Flaws — All Caught by Manual Review")

    # ── TOP: AI flaw table (the primary evidence) ──────────────────────────
    section_label(slide, "Evidence: AI Audit Findings (CI/CD YAML + Dockerfile)",
                  Inches(0.3), Inches(1.15), Inches(9.0))

    headers = ["Flaw Found in AI Output", "Risk if Shipped", "Fix Applied"]
    rows = [
        ["Hard-coded JWT_SECRET in YAML\n(e.g. JWT_SECRET: \"supersecret123\")",
         "Credentials permanently visible in git history to anyone with repo access",
         "${{ secrets.JWT_SECRET || 'safe-ci-fallback' }}\nreferenced from GitHub Secrets"],
        ["continue-on-error: true\non the Jest test step",
         "Failing tests produce a green pipeline;\nbroken code gets packaged into Docker",
         "Flag removed entirely — test failure\nnow kills the pipeline (quality gate)"],
        ["No USER directive\nin Dockerfile runtime stage",
         "Container process runs as root;\nRCE exploit → attacker gets host root",
         "addgroup nodejs + adduser nextjs\nUSER nextjs (UID 1001) before CMD"],
    ]
    add_table(slide, headers, rows,
              Inches(0.3), Inches(1.52), Inches(12.73), Inches(3.3),
              col_widths=[Inches(3.2), Inches(4.73), Inches(4.8)])

    # ── BOTTOM: Technical challenges + lesson ──────────────────────────────
    section_label(slide, "5 Technical Challenges Also Resolved",
                  Inches(0.3), Inches(4.95), Inches(9.0))

    challenges = [
        ("ESM/CJS interop  →  Next.js webpack layer", 11, False, DGRAY),
        ("Prometheus singleton lost on hot reload  →  global.__campusMetrics pattern", 11, False, DGRAY),
        ("sync(alter:true) wiped all data on restart  →  changed to sync()", 11, False, DGRAY),
        ("Middleware secret mismatch → auth loop  →  aligned 3-level fallback chains", 11, False, DGRAY),
        ("OAuth users got 401 on every API call  →  /api/auth/token bridge + JwtSynchronizer", 11, False, DGRAY),
    ]
    multiline_box(slide, challenges, Inches(0.3), Inches(5.3), Inches(9.5), Inches(2.0))

    # ── BOTTOM RIGHT: Lesson ───────────────────────────────────────────────
    lesson = [
        ("Lesson:", 13, True, NAVY),
        ("", 4),
        ("AI optimizes for working code,", 13, False, DGRAY),
        ("not security defaults.", 13, False, DGRAY),
        ("", 4),
        ("Every AI-generated DevOps", 13, False, DGRAY),
        ("artifact needs a manual", 13, False, DGRAY),
        ("security review before commit.", 13, False, DGRAY),
    ]
    multiline_box(slide, lesson, Inches(10.0), Inches(4.95), Inches(3.0), Inches(2.4))
    return slide


# ---------------------------------------------------------------------------
# Slide 8 – Testing
# Assertion: 320 Tests Across Three Isolation Layers Form a CI Quality Gate
# Evidence: Three-column breakdown + quality gate rule
# ---------------------------------------------------------------------------

def build_slide8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "320 Tests Across Three Isolation Layers Form a CI Quality Gate")

    col_w = Inches(4.15)
    gap   = Inches(0.13)

    titles = [
        "Unit Tests — 196 cases, 7 files",
        "Frontend Tests — 124 cases, 7 files",
        "Integration + CI Enforcement",
    ]
    columns = [
        [
            ("Fast, isolated — no external dependencies", 12, True, NAVY),
            ("", 4),
            ("auth.test.js", 12, False, DGRAY),
            ("  JWT sign/verify, expiry edge cases", 11, False, DGRAY),
            ("validation.test.js", 12, False, DGRAY),
            ("  All field constraints, enum values", 11, False, DGRAY),
            ("rateLimit.test.js", 12, False, DGRAY),
            ("  Window reset, IP extraction, limits", 11, False, DGRAY),
            ("eventEmitter.test.js", 12, False, DGRAY),
            ("  Domain event handlers, async flow", 11, False, DGRAY),
            ("apiErrors.test.js", 12, False, DGRAY),
            ("  Error class hierarchy, HTTP codes", 11, False, DGRAY),
            ("weatherService.test.js", 12, False, DGRAY),
            ("  Cache hits, coordinate validation", 11, False, DGRAY),
            ("middleware.test.js", 12, False, DGRAY),
            ("  Security headers, auth redirect", 11, False, DGRAY),
        ],
        [
            ("React Testing Library + jsdom — no browser", 12, True, NAVY),
            ("", 4),
            ("LoginPage.test.js", 12, False, DGRAY),
            ("  Form submit, errors, OAuth button", 11, False, DGRAY),
            ("MetricForm.test.js", 12, False, DGRAY),
            ("  Validation messages, edit mode", 11, False, DGRAY),
            ("Navbar.test.js", 12, False, DGRAY),
            ("  Auth-aware link rendering", 11, False, DGRAY),
            ("WeatherWidget.test.js", 12, False, DGRAY),
            ("  Fetch, loading, error states", 11, False, DGRAY),
            ("AuthProvider.test.js", 12, False, DGRAY),
            ("  NextAuth session wrapper", 11, False, DGRAY),
            ("apiClient.test.js", 12, False, DGRAY),
            ("  Token injection, 401 redirect", 11, False, DGRAY),
            ("Spinner.test.js", 12, False, DGRAY),
            ("  Loading UI component", 11, False, DGRAY),
        ],
        [
            ("Real HTTP calls against live server :3001", 12, True, NAVY),
            ("", 4),
            ("api.test.js  (Node --test runner)", 12, False, DGRAY),
            ("  register → login", 11, False, DGRAY),
            ("  → CRUD metrics (create/read/update/delete)", 11, False, DGRAY),
            ("  → weather endpoint fetch", 11, False, DGRAY),
            ("  → domain events table", 11, False, DGRAY),
            ("  Run locally — require server on :3001", 11, False, DGRAY),
            ("", 4),
            ("CI quality gate:", 12, True, NAVY),
            ("  Jest runs on every push to main", 11, False, DGRAY),
            ("  Pipeline FAILS if any test fails", 11, False, RED_DK),
            ("  Coverage report artifact (14 days)", 11, False, DGRAY),
        ],
    ]

    for i, (title, col_data) in enumerate(zip(titles, columns)):
        left = Inches(0.2) + i * (col_w + gap)
        section_label(slide, title, left, Inches(1.15), col_w)
        multiline_box(slide, col_data, left, Inches(1.52), col_w, Inches(5.85))
        if i < 2:
            rect(slide, left + col_w + Inches(0.05),
                 Inches(1.1), Inches(0.04), Inches(6.3), STEEL)
    return slide


# ---------------------------------------------------------------------------
# Slide 9 – Reflections
# Assertion: Observability Must Be Designed In — Not Bolted On
# Evidence: Three columns — what was built, lessons, what breaks at scale
# ---------------------------------------------------------------------------

def build_slide9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide,
        "Observability Must Be Designed In — Not Bolted On")

    col_w = Inches(4.15)
    gap   = Inches(0.13)

    titles = [
        "Evidence: Built From Scratch",
        "Evidence: Lessons That Transfer to Production",
        "Evidence: What Breaks at Scale",
    ]
    columns = [
        [
            ("Zero-dependency Prometheus module", 13, True, NAVY),
            ("lib/metrics.js: 170 lines, pure Node.js", 12, False, DGRAY),
            ("No prom-client — implements text-format 0.0.4", 12, False, DGRAY),
            ("Global singleton survives hot reloads", 12, False, DGRAY),
            ("", 5),
            ("Multi-layer auth with OAuth bridge", 13, True, NAVY),
            ("JWT + NextAuth + /api/auth/token endpoint", 12, False, DGRAY),
            ("One getAuthUser() handles all three paths", 12, False, DGRAY),
            ("", 5),
            ("Week 7 QA found bugs unit tests missed:", 13, True, NAVY),
            ("sync(alter:true) silently wiped seed data", 12, False, DGRAY),
            ("Middleware secret mismatch → redirect loop", 12, False, DGRAY),
            ("Both found only via live end-to-end testing", 12, False, DGRAY),
        ],
        [
            ("Log JSON before you ever need it", 13, True, NAVY),
            ("Shared lib/logger.js made per-route logging", 12, False, DGRAY),
            ("a 2-line add — not a 15-file refactor", 12, False, DGRAY),
            ("", 5),
            ("AI drafts; humans review for security", 13, True, NAVY),
            ("3 flaws in AI-generated CI YAML", 12, False, DGRAY),
            ("2 infra bugs surfaced only in live QA", 12, False, DGRAY),
            ("Security review is non-negotiable", 12, False, DGRAY),
            ("", 5),
            ("End-to-end testing finds what unit tests miss", 13, True, NAVY),
            ("All 320 tests passed throughout development", 12, False, DGRAY),
            ("Runtime failures emerged only in the live app", 12, False, DGRAY),
            ("", 5),
            ("\"These patterns appear in every production", 12, False, NAVY),
            ("system I will work on. This project gave me", 12, False, NAVY),
            ("hands-on practice with all of them.\"", 12, False, NAVY),
        ],
        [
            ("SQLite → PostgreSQL + connection pooling", 13, True, NAVY),
            ("SQLite write-locks under concurrent load", 12, False, DGRAY),
            ("", 5),
            ("In-memory rate limiter → Redis-backed", 13, True, NAVY),
            ("State is lost on pod restart in Kubernetes", 12, False, DGRAY),
            ("", 5),
            ("Single instance → Kubernetes + HPA", 13, True, NAVY),
            ("Auto-scale on CPU / request-rate metrics", 12, False, DGRAY),
            ("", 5),
            ("Developer experience gaps to close:", 13, True, NAVY),
            ("OpenAPI/Swagger spec for all 15 routes", 12, False, DGRAY),
            ("Grafana alerts → PagerDuty / Slack", 12, False, DGRAY),
            ("", 5),
            ("AI features planned:", 13, True, NAVY),
            ("Natural language metric search", 12, False, DGRAY),
            ("Anomaly detection on time-series data", 12, False, DGRAY),
        ],
    ]

    for i, (title, col_data) in enumerate(zip(titles, columns)):
        left = Inches(0.2) + i * (col_w + gap)
        section_label(slide, title, left, Inches(1.15), col_w)
        multiline_box(slide, col_data, left, Inches(1.52), col_w, Inches(5.85))
        if i < 2:
            rect(slide, left + col_w + Inches(0.05),
                 Inches(1.1), Inches(0.04), Inches(6.3), STEEL)
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides (assertion-evidence format)...")
    build_slide1(prs);  print("  [1/9] Title")
    build_slide2(prs);  print("  [2/9] Campus Analytics Eliminates University Metrics Data Silos")
    build_slide3(prs);  print("  [3/9] Next.js 15 Collapses Frontend and API Into One Deploy Unit")
    build_slide4(prs);  print("  [4/9] Every Technology Choice Minimizes Overhead and Maximizes Testability")
    build_slide5(prs);  print("  [5/9] Multi-Layer Defense and Seven Features Deliver Production-Grade Security")
    build_slide6(prs);  print("  [6/9] Every Commit Is Automatically Built, Tested, and Packaged")
    build_slide7(prs);  print("  [7/9] AI-Generated Code Contained 3 Security Flaws — All Caught by Manual Review")
    build_slide8(prs);  print("  [8/9] 320 Tests Across Three Isolation Layers Form a CI Quality Gate")
    build_slide9(prs);  print("  [9/9] Observability Must Be Designed In — Not Bolted On")

    prs.save(OUTPUT_PATH)
    size_kb = os.path.getsize(OUTPUT_PATH) / 1024
    print(f"\nSaved: {OUTPUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
